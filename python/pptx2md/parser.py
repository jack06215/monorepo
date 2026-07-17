"""PPTX Parser."""

import base64
import functools
import operator

import pptx.enum.shapes
import pptx.parts.image
import pptx.presentation
import pptx.shapes.autoshape
import pptx.shapes.shapetree
import pptx.table

from common.logging_util import get_logger
from pptx2md import pptx_util, types

_LOGGER = get_logger(__name__)


class PptxShapesParser:
    """PPTX Parser."""

    def __init__(self, config: types.ConversionConfig) -> None:
        """Constructor."""
        self.config = config
        self.picture_count = 0

    def process_title(
        self,
        shape: pptx.shapes.autoshape.Shape,
    ) -> types.TitleElement:
        """Process pptx shape to title element."""
        text = shape.text_frame.text.strip()
        return types.TitleElement(content=text.strip(), level=1)

    def process_table(
        self,
        shape: pptx.shapes.autoshape.Shape,
    ) -> types.TableElement | None:
        """Process pptx shape to table element."""
        if not shape.has_table:
            return None
        table_pptx: pptx.table.Table = operator.attrgetter("table")(shape)
        table: list[list[list[types.TextRun]]] = [
            [
                functools.reduce(
                    operator.iadd,
                    (pptx_util.get_text_runs(p) for p in cell.text_frame.paragraphs),
                    [],
                )
                for cell in row.cells
            ]
            for row in table_pptx.rows
        ]
        if len(table) > 0:
            return types.TableElement(content=table)
        return None

    def process_text_blocks(
        self,
        shape: pptx.shapes.autoshape.Shape,
    ) -> list[types.SlideElement]:
        """Process pptx text blocks to slide elements."""
        results: list[types.SlideElement] = []
        if pptx_util.is_list_block(shape):
            for para in shape.text_frame.paragraphs:
                if para.text.strip() == "":
                    continue
                text = pptx_util.get_text_runs(para)
                results.append(
                    types.ListItemElement(
                        content=text,
                        level=para.level,
                    )
                )
        else:
            # paragraph block
            for para in shape.text_frame.paragraphs:
                if para.text.strip() == "":
                    continue
                text = pptx_util.get_text_runs(para)
                results.append(types.ParagraphElement(content=text))
        return results

    def process_picture(
        self,
        shape: pptx.shapes.autoshape.Shape,
    ) -> types.ImageElement | None:
        """Process pptx shape to image element."""
        if self.config.disable_image:
            return None

        try:
            image_pptx: pptx.parts.image.Image = operator.attrgetter("image")(shape)
            image = base64.b64encode(image_pptx.blob)
            self.picture_count += 1
            image_str = "data:{content_type};base64,{content}"
            return types.ImageElement(
                content=image_str.format(
                    content_type=image_pptx.content_type,
                    content=image.decode(),
                )
            )

        except Exception:
            _LOGGER.error("Failed to process shape, skipped: %s", shape)
            return None


def process_shapes(
    config: types.ConversionConfig,
    shapes: list[pptx.shapes.autoshape.Shape],
) -> list[types.SlideElement]:
    """Process list of pptx shapes to slide elements."""
    pptx_shape_parser = PptxShapesParser(config=config)
    results: list[types.SlideElement] = []
    for shape in shapes:
        if pptx_util.is_title(shape):
            results.append(
                pptx_shape_parser.process_title(
                    shape=shape,
                )
            )
        elif pptx_util.is_text_block(min_block_size=config.min_block_size, shape=shape):
            results.extend(
                pptx_shape_parser.process_text_blocks(
                    shape,
                )
            )
        elif shape.shape_type == pptx.enum.shapes.MSO_SHAPE_TYPE.TABLE:
            table = pptx_shape_parser.process_table(shape)
            if table:
                results.append(table)

        elif shape.shape_type == pptx.enum.shapes.MSO_SHAPE_TYPE.PICTURE:
            try:
                picture = pptx_shape_parser.process_picture(shape)
                if picture:
                    results.append(picture)
            except Exception as e:
                _LOGGER.warning("Failed to process picture, skipped: %s", e)

    return results


def parse(
    config: types.ConversionConfig,
    prs: pptx.presentation.Presentation,
) -> types.ParsedPresentation:
    """Parses a pptx file."""
    result = types.ParsedPresentation(slides=[])

    for slide in prs.slides:
        result_slide: types.GeneralSlide | None = None
        try:
            shapes = pptx_util.ungroup_shapes(slide.shapes)
        except Exception:
            _LOGGER.warning(
                "Bad shapes encountered in this slide."
                " Please check or remove them and try again."
            )
            continue

        result_slide = types.GeneralSlide(
            elements=process_shapes(
                config,
                shapes,
            )
        )
        if not config.disable_notes and slide.has_notes_slide:
            text = (
                None
                if slide.notes_slide.notes_text_frame is None
                else slide.notes_slide.notes_text_frame.text
            )
            if text:
                result_slide.notes.append(text)
        result.slides.append(result_slide)

    return result
