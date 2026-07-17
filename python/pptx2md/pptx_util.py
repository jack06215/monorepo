"""Utility functions to check pptx properties and shared operations for working with pptx-python."""  # noqa: E501

import operator
from typing import cast

import pptx.enum.dml
import pptx.enum.shapes
import pptx.shapes.autoshape
import pptx.shapes.group
import pptx.shapes.shapetree
import pptx.text.text

from common import logging_util
from pptx2md import types

_LOGGER = logging_util.get_logger(__name__)


def is_title(shape: pptx.shapes.autoshape.Shape) -> bool:
    """Check if the shape is a title."""
    if shape.is_placeholder and (
        shape.placeholder_format.type == pptx.enum.shapes.PP_PLACEHOLDER.TITLE
        or shape.placeholder_format.type == pptx.enum.shapes.PP_PLACEHOLDER.SUBTITLE
        or shape.placeholder_format.type
        == pptx.enum.shapes.PP_PLACEHOLDER.VERTICAL_TITLE
        or shape.placeholder_format.type == pptx.enum.shapes.PP_PLACEHOLDER.CENTER_TITLE
    ):
        return True
    return False


def is_text_block(
    min_block_size: int,
    shape: pptx.shapes.autoshape.Shape,
) -> bool:
    """Check if the shape is a text block."""
    if shape.has_text_frame:
        if (
            shape.is_placeholder
            and shape.placeholder_format.type == pptx.enum.shapes.PP_PLACEHOLDER.BODY
        ):
            return True
        if len(shape.text) > min_block_size:
            return True
    return False


def is_list_block(shape: pptx.shapes.autoshape.Shape) -> bool:
    """Check if the shape is a list block."""
    levels = []
    for para in shape.text_frame.paragraphs:
        if para.level not in levels:
            levels.append(para.level)
        if para.level != 0 or len(levels) > 1:
            return True
    return False


def is_accent(font: pptx.text.text.Font) -> bool:
    """Check if the text accent text."""
    if (
        font.underline
        or font.italic
        or (
            font.color.type == pptx.enum.dml.MSO_COLOR_TYPE.SCHEME
            and (
                font.color.theme_color == pptx.enum.dml.MSO_THEME_COLOR.ACCENT_1
                or font.color.theme_color == pptx.enum.dml.MSO_THEME_COLOR.ACCENT_2
                or font.color.theme_color == pptx.enum.dml.MSO_THEME_COLOR.ACCENT_3
                or font.color.theme_color == pptx.enum.dml.MSO_THEME_COLOR.ACCENT_4
                or font.color.theme_color == pptx.enum.dml.MSO_THEME_COLOR.ACCENT_5
                or font.color.theme_color == pptx.enum.dml.MSO_THEME_COLOR.ACCENT_6
            )
        )
    ):
        return True
    return False


def is_strong(font: pptx.text.text.Font) -> bool:
    """Check if the shape is bold."""
    if font.bold or (
        font.color.type == pptx.enum.dml.MSO_COLOR_TYPE.SCHEME
        and (
            font.color.theme_color == pptx.enum.dml.MSO_THEME_COLOR.DARK_1
            or font.color.theme_color == pptx.enum.dml.MSO_THEME_COLOR.DARK_2
        )
    ):
        return True
    return False


def get_text_runs(para: pptx.text.text._Paragraph) -> list[types.TextRun]:
    """Parse paragraph into style-aware text block."""
    runs = []
    for run in para.runs:
        result = types.TextRun(text=run.text, style=types.TextStyle())
        if result.text == "":
            continue
        try:
            if run.hyperlink.address:
                result.style.hyperlink = run.hyperlink.address
        except Exception:
            result.style.hyperlink = "error:ppt-link-parsing-issue"
        if is_accent(run.font):
            result.style.is_accent = True
        if is_strong(run.font):
            result.style.is_strong = True
        if run.font.color.type == pptx.enum.dml.MSO_COLOR_TYPE.RGB:
            result.style.color_rgb = run.font.color.rgb
        runs.append(result)
    return runs


def ungroup_shapes(
    shapes: pptx.shapes.shapetree.SlideShapes | pptx.shapes.shapetree.GroupShapes,
) -> list[pptx.shapes.autoshape.Shape]:
    """Deconstruct group shapes."""
    res = []
    for shape in shapes:
        try:
            if shape.shape_type == pptx.enum.shapes.MSO_SHAPE_TYPE.GROUP:
                res.extend(
                    ungroup_shapes(cast(pptx.shapes.group.GroupShape, shape).shapes)
                )
            else:
                res.append(cast(pptx.shapes.autoshape.Shape, shape))
        except Exception as e:
            _LOGGER.warning("failed to ungroup shape %s, skipped. Error: %s", shape, e)
    return sorted(res, key=operator.attrgetter("top", "left"))
